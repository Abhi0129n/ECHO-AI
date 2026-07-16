import os
import csv
import openpyxl
from docx import Document
from pptx import Presentation
from typing import List, Dict, Any, Optional
from tools.productivity.schemas import CreateExcelRequest, CreateWordRequest, CreatePPTRequest
from tools.productivity.utils import resolve_and_verify_path

class ProductivityService:
    def __init__(self, base_dir: str = "."):
        self.base_dir = os.path.abspath(base_dir)

    def create_excel(self, path: str, data: Dict[str, List[List[Any]]]) -> str:
        abs_path = resolve_and_verify_path(self.base_dir, path, [".xlsx"])
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        
        wb = openpyxl.Workbook()
        # Remove default active sheet to create sheets exactly as requested
        if wb.active:
            wb.remove(wb.active)
            
        for sheet_name, rows in data.items():
            ws = wb.create_sheet(title=sheet_name)
            for row in rows:
                ws.append(row)
                
        wb.save(abs_path)
        return os.path.relpath(abs_path, self.base_dir)

    def create_word(self, path: str, title: str, paragraphs: List[Any]) -> str:
        abs_path = resolve_and_verify_path(self.base_dir, path, [".docx"])
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        
        doc = Document()
        doc.add_heading(title, level=0)
        
        for p in paragraphs:
            if p.is_heading:
                doc.add_heading(p.text, level=p.heading_level)
            else:
                doc.add_paragraph(p.text)
                
        doc.save(abs_path)
        return os.path.relpath(abs_path, self.base_dir)

    def create_ppt(self, path: str, title: str, subtitle: Optional[str], slides: List[Any]) -> str:
        abs_path = resolve_and_verify_path(self.base_dir, path, [".pptx"])
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        
        prs = Presentation()
        
        # Add title slide (layout index 0)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
            
        # Add content slides (layout index 1 is Title and Content)
        content_layout = prs.slide_layouts[1]
        for s in slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = s.title
            slide.placeholders[1].text = s.content
            
        prs.save(abs_path)
        return os.path.relpath(abs_path, self.base_dir)

    def read_csv(self, path: str) -> List[List[str]]:
        abs_path = resolve_and_verify_path(self.base_dir, path, [".csv"])
        if not os.path.exists(abs_path):
            raise FileNotFoundError(f"CSV file '{path}' does not exist")
            
        rows = []
        with open(abs_path, mode="r", newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        return rows
